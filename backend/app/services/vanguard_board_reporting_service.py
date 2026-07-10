"""v4.6 — Project Vanguard, Section 7: Board Reporting.

Extends `atlas_report_service.generate_executive_report` (already
audience/cadence-typed and persisted) rather than the three other board-
report generators already in this codebase — `routes/board_reporting.py`
(single fixed shape, no audience typing), `benchmark_engine.
generate_board_report` (a second, CVInferenceRecord-based enterprise
rollup lineage), and `portfolio_briefings.py` (a different audience
entirely — LumenAI's own SaaS customer portfolio, not the hospital's
board). None of those are read from or extended here.

Four named packet types map onto Atlas's existing `audience`/`cadence`
enum (never a new one): a tenant with no resolvable enterprise-hierarchy
facility falls back to the Vanguard Executive Intelligence Center
snapshot instead of a fabricated Atlas-shaped report.

Export formats: PDF (the same `reportlab` low-level canvas pattern
`atlas_report_service.build_report_pdf_bytes` already uses), Excel (the
same `openpyxl` pattern), and PowerPoint — genuinely new for a hospital-
facing report, but built with the same `python-pptx` `Presentation()`
pattern already used internally by `leadership_packet_exports.py`/
`governance_packet_exports.py`, not a new export library.
"""
from __future__ import annotations

import json
from io import BytesIO

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy.orm import Session

from app.models.vanguard_intelligence import (
    BOARD_PACKET_TYPES,
    DISCLAIMER,
    PACKET_ANNUAL_STRATEGIC,
    PACKET_MONTHLY_BOARD,
    PACKET_QUALITY_COMMITTEE,
    PACKET_QUARTERLY_REVIEW,
    BoardReportPacket,
)
from app.services import atlas_report_service, finding_trend_service, platform_org_service, vanguard_executive_intelligence_service

_PACKET_AUDIENCE_CADENCE = {
    PACKET_MONTHLY_BOARD: ("ceo", "monthly"),
    PACKET_QUARTERLY_REVIEW: ("coo", "quarterly"),
    PACKET_ANNUAL_STRATEGIC: ("ceo", "annual"),
    PACKET_QUALITY_COMMITTEE: ("spd_director", "monthly"),
}

_PACKET_TITLES = {
    PACKET_MONTHLY_BOARD: "Monthly Board Packet",
    PACKET_QUARTERLY_REVIEW: "Quarterly Executive Review",
    PACKET_ANNUAL_STRATEGIC: "Annual Strategic Report",
    PACKET_QUALITY_COMMITTEE: "Quality Committee Report",
}


class UnknownPacketTypeError(Exception):
    pass


class PacketNotFoundError(Exception):
    pass


def _row_to_dict(row: BoardReportPacket) -> dict:
    return {
        "id": row.id, "packet_type": row.packet_type, "title": row.title, "content": json.loads(row.content_json),
        "generated_by": row.generated_by, "created_at": row.created_at.isoformat(),
        "human_review_required": row.human_review_required, "disclaimer": row.disclaimer,
    }


def generate_board_packet(db: Session, tenant_id: str, packet_type: str, *, generated_by: str = "system") -> dict:
    if packet_type not in BOARD_PACKET_TYPES:
        raise UnknownPacketTypeError(f"packet_type must be one of {BOARD_PACKET_TYPES}")

    audience, cadence = _PACKET_AUDIENCE_CADENCE[packet_type]
    executive_snapshot = vanguard_executive_intelligence_service.executive_intelligence_center(db, tenant_id)

    facility = platform_org_service.facility_for_tenant(db, tenant_id)
    if facility is not None:
        atlas_report = atlas_report_service.generate_executive_report(
            db, facility["system_id"], audience=audience, cadence=cadence, generated_by=generated_by,
        )
        content = {"source": "atlas_report", "atlas_report": atlas_report, "executive_snapshot": executive_snapshot}
    else:
        content = {
            "source": "vanguard_executive_snapshot", "executive_snapshot": executive_snapshot,
            "note": "No enterprise-hierarchy facility on record — using the live Executive Intelligence Center snapshot rather than a fabricated Atlas-shaped report.",
        }

    if packet_type == PACKET_QUALITY_COMMITTEE:
        content["finding_trends"] = finding_trend_service.finding_trends(db, tenant_id)

    row = BoardReportPacket(
        tenant_id=tenant_id, packet_type=packet_type, title=_PACKET_TITLES[packet_type],
        content_json=json.dumps(content, default=str), generated_by=generated_by,
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return _row_to_dict(row)


def get_packet(db: Session, tenant_id: str, packet_id: int) -> dict:
    row = db.query(BoardReportPacket).filter(BoardReportPacket.id == packet_id, BoardReportPacket.tenant_id == tenant_id).first()
    if row is None:
        raise PacketNotFoundError(f"Board report packet {packet_id} not found for tenant {tenant_id}.")
    return _row_to_dict(row)


def list_packets(db: Session, tenant_id: str, *, packet_type: str = "") -> list[dict]:
    q = db.query(BoardReportPacket).filter(BoardReportPacket.tenant_id == tenant_id)
    if packet_type:
        q = q.filter(BoardReportPacket.packet_type == packet_type)
    return [_row_to_dict(r) for r in q.order_by(BoardReportPacket.id.desc()).all()]


def _flatten_metrics(content: dict) -> list[tuple[str, str]]:
    snapshot = content.get("executive_snapshot", {})
    rows = [
        ("Enterprise Risk Score", snapshot.get("enterprise_readiness", {}).get("enterprise_risk_score")),
        ("Surgical Readiness %", snapshot.get("surgical_readiness", {}).get("readiness_pct")),
        ("Repair Cost Trend (USD)", snapshot.get("financial_impact", {}).get("repair_cost_trend_usd")),
        ("Avoided Replacement Cost (USD)", snapshot.get("financial_impact", {}).get("avoided_replacement_cost_usd")),
        ("Capacity Utilization %", snapshot.get("capacity", {}).get("utilization_pct")),
    ]
    return [(k, str(v)) for k, v in rows]


def build_packet_pdf_bytes(packet: dict) -> bytes:
    bio = BytesIO()
    c = canvas.Canvas(bio, pagesize=letter)
    width, height = letter
    y = height - 50

    c.setFont("Helvetica-Bold", 18)
    c.drawString(50, y, packet["title"])
    y -= 26
    c.setFont("Helvetica", 10)
    c.drawString(50, y, f"Generated: {packet['created_at']} · By: {packet['generated_by']}")
    y -= 24

    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Executive Summary")
    y -= 18
    c.setFont("Helvetica", 10)
    for key, value in _flatten_metrics(packet["content"]):
        c.drawString(60, y, f"{key}: {value}")
        y -= 14

    y -= 10
    c.setFont("Helvetica-Oblique", 8)
    for line in [DISCLAIMER[i:i + 100] for i in range(0, len(DISCLAIMER), 100)]:
        c.drawString(50, y, line)
        y -= 10

    c.showPage()
    c.save()
    return bio.getvalue()


def build_packet_xlsx_bytes(packet: dict) -> bytes:
    wb = Workbook()
    summary_ws = wb.active
    summary_ws.title = "Summary"
    summary_ws.append(["Metric", "Value"])
    for key, value in _flatten_metrics(packet["content"]):
        summary_ws.append([key, value])

    detail_ws = wb.create_sheet("Detail")
    detail_ws.append(["Key", "Value"])
    for key, value in packet["content"].items():
        if isinstance(value, (dict, list)):
            detail_ws.append([key, json.dumps(value, default=str)[:500]])
        else:
            detail_ws.append([key, value])

    bio = BytesIO()
    wb.save(bio)
    return bio.getvalue()


def build_packet_pptx_bytes(packet: dict) -> bytes:
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = packet["title"]
    slide.placeholders[1].text = f"Generated {packet['created_at']} · {packet['generated_by']}"

    metrics_layout = prs.slide_layouts[1]
    metrics_slide = prs.slides.add_slide(metrics_layout)
    metrics_slide.shapes.title.text = "Executive Summary"
    body = metrics_slide.placeholders[1].text_frame
    body.clear()
    metrics = _flatten_metrics(packet["content"])
    if metrics:
        body.text = f"{metrics[0][0]}: {metrics[0][1]}"
        for key, value in metrics[1:]:
            p = body.add_paragraph()
            p.text = f"{key}: {value}"
            p.font.size = Pt(16)

    disclaimer_slide = prs.slides.add_slide(prs.slide_layouts[1])
    disclaimer_slide.shapes.title.text = "Disclaimer"
    disclaimer_slide.placeholders[1].text_frame.text = packet["disclaimer"]

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
