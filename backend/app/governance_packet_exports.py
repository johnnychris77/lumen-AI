from __future__ import annotations

import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy import text
from sqlalchemy.orm import Session

from app.executive_escalations import generate_governance_packet


ARTIFACT_ROOT = Path("generated_governance_packets")


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def ensure_governance_packet_tables(db: Session) -> None:
    db.execute(
        text(
            """
            CREATE TABLE IF NOT EXISTS executive_governance_packets (
                id SERIAL PRIMARY KEY,
                packet_title VARCHAR(255) NOT NULL,
                executive_summary TEXT NOT NULL DEFAULT '',
                packet_json TEXT NOT NULL DEFAULT '{}',
                status VARCHAR(50) NOT NULL DEFAULT 'generated',
                created_at TIMESTAMP WITH TIME ZONE NOT NULL DEFAULT NOW()
            )
            """
        )
    )

    db.execute(
        text(
            """
            CREATE TABLE IF NOT EXISTS executive_governance_packet_exports (
                id SERIAL PRIMARY KEY,
                packet_id INTEGER NOT NULL,
                export_title VARCHAR(255) NOT NULL,
                format_bundle VARCHAR(100) NOT NULL DEFAULT 'docx_pptx_pdf',
                docx_path TEXT,
                pptx_path TEXT,
                pdf_path TEXT,
                created_at TIMESTAMP WITH TIME ZONE NOT NULL DEFAULT NOW()
            )
            """
        )
    )

    db.execute(
        text(
            """
            CREATE TABLE IF NOT EXISTS executive_governance_packet_deliveries (
                id SERIAL PRIMARY KEY,
                packet_id INTEGER NOT NULL,
                export_id INTEGER,
                delivery_channel VARCHAR(50) NOT NULL DEFAULT 'internal',
                delivery_target TEXT NOT NULL DEFAULT 'executive-governance-council',
                status VARCHAR(50) NOT NULL DEFAULT 'sent',
                payload_json TEXT NOT NULL DEFAULT '{}',
                created_at TIMESTAMP WITH TIME ZONE NOT NULL DEFAULT NOW()
            )
            """
        )
    )

    db.commit()


def create_governance_packet_record(db: Session, packet_title: str | None = None) -> dict[str, Any]:
    ensure_governance_packet_tables(db)

    packet = generate_governance_packet(db)
    title = packet_title or f"Executive Governance Packet — {_now_iso()}"

    row = (
        db.execute(
            text(
                """
                INSERT INTO executive_governance_packets (
                    packet_title,
                    executive_summary,
                    packet_json,
                    status
                )
                VALUES (
                    :packet_title,
                    :executive_summary,
                    :packet_json,
                    'generated'
                )
                RETURNING *
                """
            ),
            {
                "packet_title": title,
                "executive_summary": packet.get("executive_summary", ""),
                "packet_json": json.dumps(packet, default=str),
            },
        )
        .mappings()
        .first()
    )

    db.commit()
    return dict(row)


def get_governance_packet(db: Session, packet_id: int) -> dict[str, Any] | None:
    ensure_governance_packet_tables(db)

    row = (
        db.execute(
            text(
                """
                SELECT *
                FROM executive_governance_packets
                WHERE id = :packet_id
                """
            ),
            {"packet_id": packet_id},
        )
        .mappings()
        .first()
    )

    return dict(row) if row else None


def list_governance_packets(db: Session, limit: int = 50) -> list[dict[str, Any]]:
    ensure_governance_packet_tables(db)

    rows = (
        db.execute(
            text(
                """
                SELECT *
                FROM executive_governance_packets
                ORDER BY created_at DESC, id DESC
                LIMIT :limit
                """
            ),
            {"limit": limit},
        )
        .mappings()
        .all()
    )

    return [dict(row) for row in rows]


def _packet_payload(packet_record: dict[str, Any]) -> dict[str, Any]:
    try:
        return json.loads(packet_record.get("packet_json") or "{}")
    except Exception:
        return {}


def _write_docx(packet_record: dict[str, Any], output_path: Path) -> None:
    payload = _packet_payload(packet_record)

    doc = Document()
    doc.add_heading(packet_record["packet_title"], 0)
    doc.add_paragraph(packet_record.get("executive_summary") or "")

    doc.add_heading("Governance Rollup", level=1)
    rollup = payload.get("rollup", {})
    for key, value in rollup.items():
        doc.add_paragraph(f"{key}: {value}")

    doc.add_heading("Top Governance Items", level=1)
    for item in payload.get("top_governance_items", []):
        doc.add_heading(str(item.get("tenant") or "Tenant"), level=2)
        doc.add_paragraph(f"Priority: {item.get('priority')}")
        doc.add_paragraph(f"Type: {item.get('type')}")
        doc.add_paragraph(f"Owner: {item.get('owner')}")
        doc.add_paragraph(str(item.get("summary") or ""))

    doc.add_heading("Recommended Leadership Decisions", level=1)
    for decision in payload.get("recommended_leadership_decisions", []):
        doc.add_paragraph(str(decision), style="List Bullet")

    doc.save(output_path)


def _write_pptx(packet_record: dict[str, Any], output_path: Path) -> None:
    payload = _packet_payload(packet_record)
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = packet_record["packet_title"]
    slide.placeholders[1].text = packet_record.get("executive_summary") or ""

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Governance Rollup"
    rollup = payload.get("rollup", {})
    slide.placeholders[1].text = "\n".join([f"{k}: {v}" for k, v in rollup.items()])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Governance Items"
    items = payload.get("top_governance_items", [])[:6]
    slide.placeholders[1].text = "\n".join(
        [
            f"{item.get('tenant')} — {item.get('priority')} — {item.get('owner')}"
            for item in items
        ]
    ) or "No open governance items."

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommended Leadership Decisions"
    slide.placeholders[1].text = "\n".join(payload.get("recommended_leadership_decisions", []))

    prs.save(output_path)


def _write_pdf(packet_record: dict[str, Any], output_path: Path) -> None:
    payload = _packet_payload(packet_record)
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter

    y = height - 50
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, y, packet_record["packet_title"][:80])
    y -= 30

    c.setFont("Helvetica", 10)

    def line(text_value: str) -> None:
        nonlocal y
        if y < 60:
            c.showPage()
            y = height - 50
            c.setFont("Helvetica", 10)
        c.drawString(50, y, text_value[:110])
        y -= 16

    line(packet_record.get("executive_summary") or "")
    y -= 10

    c.setFont("Helvetica-Bold", 12)
    line("Governance Rollup")
    c.setFont("Helvetica", 10)

    for key, value in (payload.get("rollup") or {}).items():
        line(f"{key}: {value}")

    y -= 10
    c.setFont("Helvetica-Bold", 12)
    line("Top Governance Items")
    c.setFont("Helvetica", 10)

    for item in payload.get("top_governance_items", []):
        line(f"{item.get('tenant')} | {item.get('priority')} | {item.get('owner')}")
        line(str(item.get("summary") or ""))

    y -= 10
    c.setFont("Helvetica-Bold", 12)
    line("Recommended Leadership Decisions")
    c.setFont("Helvetica", 10)

    for decision in payload.get("recommended_leadership_decisions", []):
        line(f"- {decision}")

    c.save()


def export_governance_packet(db: Session, packet_id: int) -> dict[str, Any]:
    ensure_governance_packet_tables(db)

    packet = get_governance_packet(db, packet_id)
    if not packet:
        raise ValueError(f"Governance packet {packet_id} was not found")

    packet_dir = ARTIFACT_ROOT / f"packet_{packet_id}"
    packet_dir.mkdir(parents=True, exist_ok=True)

    docx_path = packet_dir / f"governance_packet_{packet_id}.docx"
    pptx_path = packet_dir / f"governance_packet_{packet_id}.pptx"
    pdf_path = packet_dir / f"governance_packet_{packet_id}.pdf"

    _write_docx(packet, docx_path)
    _write_pptx(packet, pptx_path)
    _write_pdf(packet, pdf_path)

    row = (
        db.execute(
            text(
                """
                INSERT INTO executive_governance_packet_exports (
                    packet_id,
                    export_title,
                    docx_path,
                    pptx_path,
                    pdf_path
                )
                VALUES (
                    :packet_id,
                    :export_title,
                    :docx_path,
                    :pptx_path,
                    :pdf_path
                )
                RETURNING *
                """
            ),
            {
                "packet_id": packet_id,
                "export_title": f"Governance Packet Export — {packet.get('packet_title')}",
                "docx_path": str(docx_path),
                "pptx_path": str(pptx_path),
                "pdf_path": str(pdf_path),
            },
        )
        .mappings()
        .first()
    )

    db.commit()
    return dict(row)


def list_governance_packet_exports(db: Session, packet_id: int) -> list[dict[str, Any]]:
    ensure_governance_packet_tables(db)

    rows = (
        db.execute(
            text(
                """
                SELECT *
                FROM executive_governance_packet_exports
                WHERE packet_id = :packet_id
                ORDER BY created_at DESC, id DESC
                """
            ),
            {"packet_id": packet_id},
        )
        .mappings()
        .all()
    )

    return [dict(row) for row in rows]


def get_governance_packet_export(db: Session, export_id: int) -> dict[str, Any] | None:
    ensure_governance_packet_tables(db)

    row = (
        db.execute(
            text(
                """
                SELECT *
                FROM executive_governance_packet_exports
                WHERE id = :export_id
                """
            ),
            {"export_id": export_id},
        )
        .mappings()
        .first()
    )

    return dict(row) if row else None


def deliver_governance_packet(
    db: Session,
    packet_id: int,
    export_id: int | None,
    delivery_channel: str = "internal",
    delivery_target: str = "executive-governance-council",
    message: str = "Executive governance packet is ready for review.",
) -> dict[str, Any]:
    ensure_governance_packet_tables(db)

    packet = get_governance_packet(db, packet_id)
    if not packet:
        raise ValueError(f"Governance packet {packet_id} was not found")

    payload = {
        "packet_id": packet_id,
        "export_id": export_id,
        "delivery_channel": delivery_channel,
        "delivery_target": delivery_target,
        "message": message,
        "sent_at": _now_iso(),
        "note": "Development-mode governance packet delivery audit.",
    }

    row = (
        db.execute(
            text(
                """
                INSERT INTO executive_governance_packet_deliveries (
                    packet_id,
                    export_id,
                    delivery_channel,
                    delivery_target,
                    status,
                    payload_json
                )
                VALUES (
                    :packet_id,
                    :export_id,
                    :delivery_channel,
                    :delivery_target,
                    'sent',
                    :payload_json
                )
                RETURNING *
                """
            ),
            {
                "packet_id": packet_id,
                "export_id": export_id,
                "delivery_channel": delivery_channel,
                "delivery_target": delivery_target,
                "payload_json": json.dumps(payload, default=str),
            },
        )
        .mappings()
        .first()
    )

    db.commit()
    return dict(row)


def list_governance_packet_deliveries(db: Session, packet_id: int) -> list[dict[str, Any]]:
    ensure_governance_packet_tables(db)

    rows = (
        db.execute(
            text(
                """
                SELECT *
                FROM executive_governance_packet_deliveries
                WHERE packet_id = :packet_id
                ORDER BY created_at DESC, id DESC
                """
            ),
            {"packet_id": packet_id},
        )
        .mappings()
        .all()
    )

    return [dict(row) for row in rows]
