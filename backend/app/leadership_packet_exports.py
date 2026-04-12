from __future__ import annotations

import json
import os
from pathlib import Path

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy.orm import Session

from app.branding import get_branding
from app.db import models

EXPORT_DIR = Path("/app/backend/generated_packets")


def _ensure_dir() -> Path:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    return EXPORT_DIR


def _safe_prefix(value: str) -> str:
    return "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in value).strip("_") or "packet"


def _load_briefing(row: models.GeneratedBriefing) -> tuple[list[dict], dict]:
    slides = json.loads(row.slide_outline_json or "[]")
    summary = json.loads(row.summary_json or "{}")
    return slides, summary


def generate_docx_packet(
    *,
    output_path: Path,
    branding: dict,
    briefing: models.GeneratedBriefing,
) -> None:
    doc = Document()
    title = doc.add_heading(briefing.title, level=0)
    title.runs[0].font.size = Pt(20)

    p = doc.add_paragraph()
    p.add_run(f"Organization: {branding['display_name']}\n").bold = True
    p.add_run(f"Audience: {briefing.audience}\n")
    p.add_run(f"Period: {briefing.period_label}\n")
    if branding.get("support_email"):
        p.add_run(f"Support: {branding['support_email']}\n")

    if branding.get("welcome_text"):
        doc.add_paragraph(branding["welcome_text"])

    doc.add_heading("Executive Memo", level=1)
    doc.add_paragraph(briefing.memo_text)

    doc.add_heading("Slide Outline", level=1)
    slides, _ = _load_briefing(briefing)
    for idx, slide in enumerate(slides, start=1):
        doc.add_paragraph(f"Slide {idx}: {slide.get('title', '')}", style="List Bullet")
        subtitle = slide.get("subtitle", "")
        if subtitle:
            doc.add_paragraph(f"Subtitle: {subtitle}")
        for bullet in slide.get("bullets", []):
            doc.add_paragraph(str(bullet), style="List Bullet 2")

    doc.save(str(output_path))


def generate_pptx_packet(
    *,
    output_path: Path,
    branding: dict,
    briefing: models.GeneratedBriefing,
) -> None:
    prs = Presentation()
    slides, _ = _load_briefing(briefing)

    for idx, slide_data in enumerate(slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data.get("title", briefing.title)

        body = slide.placeholders[1].text_frame
        body.clear()

        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            p = body.paragraphs[0]
            p.text = subtitle
            p.level = 0

        bullets = slide_data.get("bullets", [])
        for i, bullet in enumerate(bullets):
            p = body.add_paragraph() if (subtitle or i > 0) else body.paragraphs[0]
            p.text = str(bullet)
            p.level = 0

        tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.7), Inches(6), Inches(0.3))
        tf = tx_box.text_frame
        tf.text = branding.get("display_name", briefing.tenant_name)

    prs.save(str(output_path))


def generate_pdf_packet(
    *,
    output_path: Path,
    branding: dict,
    briefing: models.GeneratedBriefing,
) -> None:
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    y = height - 50

    def write_line(text: str, font="Helvetica", size=11, advance=16):
        nonlocal y
        c.setFont(font, size)
        c.drawString(50, y, text[:110])
        y -= advance
        if y < 60:
            c.showPage()
            y = height - 50

    write_line(briefing.title, "Helvetica-Bold", 16, 22)
    write_line(f"Organization: {branding['display_name']}", "Helvetica", 11)
    write_line(f"Audience: {briefing.audience}", "Helvetica", 11)
    write_line(f"Period: {briefing.period_label}", "Helvetica", 11)
    if branding.get("support_email"):
        write_line(f"Support: {branding['support_email']}", "Helvetica", 11)
    y -= 8

    write_line("Executive Memo", "Helvetica-Bold", 13, 20)
    for line in briefing.memo_text.split("\n"):
        if line.strip():
            write_line(line, "Helvetica", 10, 14)
        else:
            y -= 8

    y -= 8
    write_line("Slide Outline", "Helvetica-Bold", 13, 20)
    slides, _ = _load_briefing(briefing)
    for idx, slide in enumerate(slides, start=1):
        write_line(f"Slide {idx}: {slide.get('title', '')}", "Helvetica-Bold", 11, 16)
        subtitle = slide.get("subtitle", "")
        if subtitle:
            write_line(f"Subtitle: {subtitle}", "Helvetica", 10, 14)
        for bullet in slide.get("bullets", []):
            write_line(f"- {bullet}", "Helvetica", 10, 14)

    c.save()


def build_leadership_packet(
    db: Session,
    *,
    tenant_id: str,
    tenant_name: str,
    briefing_id: int,
) -> models.LeadershipPacket:
    briefing = (
        db.query(models.GeneratedBriefing)
        .filter(
            models.GeneratedBriefing.id == briefing_id,
            models.GeneratedBriefing.tenant_id == tenant_id,
        )
        .first()
    )
    if not briefing:
        raise ValueError("Briefing not found")

    branding = get_branding(db, tenant_id, tenant_name)
    outdir = _ensure_dir()
    prefix = _safe_prefix(branding.get("export_prefix") or tenant_id)
    stem = f"{prefix}_briefing_{briefing.id}"

    docx_path = outdir / f"{stem}.docx"
    pptx_path = outdir / f"{stem}.pptx"
    pdf_path = outdir / f"{stem}.pdf"

    generate_docx_packet(output_path=docx_path, branding=branding, briefing=briefing)
    generate_pptx_packet(output_path=pptx_path, branding=branding, briefing=briefing)
    generate_pdf_packet(output_path=pdf_path, branding=branding, briefing=briefing)

    row = models.LeadershipPacket(
        tenant_id=tenant_id,
        tenant_name=tenant_name,
        briefing_id=briefing.id,
        packet_type=briefing.briefing_type,
        title=briefing.title,
        docx_path=str(docx_path),
        pptx_path=str(pptx_path),
        pdf_path=str(pdf_path),
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return row
