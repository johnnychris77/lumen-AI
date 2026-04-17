from __future__ import annotations

import json
from pathlib import Path

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy.orm import Session

from app.branding import get_branding
from app.db import models

EXPORT_DIR = Path("/app/backend/generated_account_reviews")


def _ensure_dir() -> Path:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    return EXPORT_DIR


def _safe_prefix(value: str) -> str:
    return "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in value).strip("_") or "account_review"


def _safe_json_obj(value: str | None) -> dict:
    if not value:
        return {}
    try:
        parsed = json.loads(value)
        return parsed if isinstance(parsed, dict) else {}
    except Exception:
        return {}


def _safe_json_list(value: str | None) -> list:
    if not value:
        return []
    try:
        parsed = json.loads(value)
        return parsed if isinstance(parsed, list) else []
    except Exception:
        return []


def generate_docx_export(*, output_path: Path, branding: dict, review: models.AccountReviewPacket) -> None:
    doc = Document()
    heading = doc.add_heading(review.title, level=0)
    heading.runs[0].font.size = Pt(20)

    p = doc.add_paragraph()
    p.add_run(f"Organization: {branding['display_name']}\n").bold = True
    p.add_run(f"Review Type: {review.review_type}\n")
    p.add_run(f"Period: {review.period_label}\n")
    if branding.get("support_email"):
        p.add_run(f"Support: {branding['support_email']}\n")

    if branding.get("welcome_text"):
        doc.add_paragraph(branding["welcome_text"])

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(review.executive_summary)

    doc.add_heading("QBR Narrative", level=1)
    doc.add_paragraph(review.qbr_narrative)

    risks = _safe_json_list(review.risks_json)
    next_steps = _safe_json_list(review.next_steps_json)
    summary = _safe_json_obj(review.summary_json)

    doc.add_heading("Key Risks", level=1)
    if risks:
        for item in risks:
            doc.add_paragraph(str(item), style="List Bullet")
    else:
        doc.add_paragraph("No material risks recorded.")

    doc.add_heading("Next Steps", level=1)
    if next_steps:
        for item in next_steps:
            doc.add_paragraph(str(item), style="List Bullet")
    else:
        doc.add_paragraph("No next steps recorded.")

    doc.add_heading("Summary Metrics", level=1)
    for k, v in summary.items():
        doc.add_paragraph(f"{k}: {v}")

    doc.save(str(output_path))


def generate_pptx_export(*, output_path: Path, branding: dict, review: models.AccountReviewPacket) -> None:
    prs = Presentation()

    slides_data = [
        {
            "title": review.title,
            "subtitle": review.period_label,
            "bullets": [
                branding.get("welcome_text", f"Executive account review for {branding['display_name']}"),
                f"Review type: {review.review_type}",
            ],
        },
        {
            "title": "Executive Summary",
            "subtitle": "",
            "bullets": [line for line in review.executive_summary.split("\n") if line.strip()][:6] or ["No executive summary available."],
        },
        {
            "title": "QBR Narrative",
            "subtitle": "",
            "bullets": [line for line in review.qbr_narrative.split("\n") if line.strip()][:6] or ["No QBR narrative available."],
        },
        {
            "title": "Key Risks",
            "subtitle": "",
            "bullets": _safe_json_list(review.risks_json)[:6] or ["No material risks recorded."],
        },
        {
            "title": "Next Steps",
            "subtitle": "",
            "bullets": _safe_json_list(review.next_steps_json)[:6] or ["No next steps recorded."],
        },
    ]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]
        body = slide.placeholders[1].text_frame
        body.clear()

        first_written = False
        if slide_data.get("subtitle"):
            p = body.paragraphs[0]
            p.text = slide_data["subtitle"]
            p.level = 0
            first_written = True

        for bullet in slide_data.get("bullets", []):
            if not first_written:
                p = body.paragraphs[0]
                first_written = True
            else:
                p = body.add_paragraph()
            p.text = str(bullet)
            p.level = 0

        tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.7), Inches(6), Inches(0.3))
        tf = tx_box.text_frame
        tf.text = branding.get("display_name", review.tenant_name)

    prs.save(str(output_path))


def generate_pdf_export(*, output_path: Path, branding: dict, review: models.AccountReviewPacket) -> None:
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    y = height - 50

    def write_line(text: str, font="Helvetica", size=11, advance=16):
        nonlocal y
        c.setFont(font, size)
        c.drawString(50, y, text[:110])
        y -= advance
        if y < 60:
            c.showPage()
            y = height - 50

    write_line(review.title, "Helvetica-Bold", 16, 22)
    write_line(f"Organization: {branding['display_name']}")
    write_line(f"Review Type: {review.review_type}")
    write_line(f"Period: {review.period_label}")
    if branding.get("support_email"):
        write_line(f"Support: {branding['support_email']}")
    y -= 8

    write_line("Executive Summary", "Helvetica-Bold", 13, 20)
    for line in review.executive_summary.split("\n"):
        if line.strip():
            write_line(line, "Helvetica", 10, 14)

    y -= 8
    write_line("QBR Narrative", "Helvetica-Bold", 13, 20)
    for line in review.qbr_narrative.split("\n"):
        if line.strip():
            write_line(line, "Helvetica", 10, 14)

    y -= 8
    write_line("Key Risks", "Helvetica-Bold", 13, 20)
    for item in _safe_json_list(review.risks_json) or ["No material risks recorded."]:
        write_line(f"- {item}", "Helvetica", 10, 14)

    y -= 8
    write_line("Next Steps", "Helvetica-Bold", 13, 20)
    for item in _safe_json_list(review.next_steps_json) or ["No next steps recorded."]:
        write_line(f"- {item}", "Helvetica", 10, 14)

    c.save()


def build_account_review_export(db: Session, *, tenant_id: str, tenant_name: str, account_review_id: int) -> models.AccountReviewExport:
    review = (
        db.query(models.AccountReviewPacket)
        .filter(
            models.AccountReviewPacket.id == account_review_id,
            models.AccountReviewPacket.tenant_id == tenant_id,
        )
        .first()
    )
    if not review:
        raise ValueError("Account review packet not found")

    branding = get_branding(db, tenant_id, tenant_name)
    outdir = _ensure_dir()
    prefix = _safe_prefix(branding.get("export_prefix") or tenant_id)
    stem = f"{prefix}_account_review_{review.id}"

    docx_path = outdir / f"{stem}.docx"
    pptx_path = outdir / f"{stem}.pptx"
    pdf_path = outdir / f"{stem}.pdf"

    generate_docx_export(output_path=docx_path, branding=branding, review=review)
    generate_pptx_export(output_path=pptx_path, branding=branding, review=review)
    generate_pdf_export(output_path=pdf_path, branding=branding, review=review)

    row = models.AccountReviewExport(
        tenant_id=tenant_id,
        tenant_name=tenant_name,
        account_review_id=review.id,
        export_type=review.review_type,
        title=review.title,
        docx_path=str(docx_path),
        pptx_path=str(pptx_path),
        pdf_path=str(pdf_path),
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return row
